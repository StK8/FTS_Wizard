import collections.abc
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re
import win32com.client
import os


from PIL import Image
from PIL import ImageOps


DIRECTORY = os.getcwd()

IMAGE_HEIGHT = 10.5

HC_Composition_slides = []
Contamination_slides = []
Sampling_summary_slides = []


# aux function to add image border
def add_border(input_image, output_image, border):
    img = Image.open(input_image)
    if isinstance(border, int) or isinstance(border, tuple):
        bimg = ImageOps.expand(img, border=border)
    else:
        raise RuntimeError('Border is not an image or tuple')
    bimg.save(output_image)

def iter_shapes(prs):
    prev_file_number = -1
    current_file_number = -1
    prev_station_depth = -1
    current_station_depth = -1


    for slide in prs.slides:
        # print(f"Slide number: {prs.slides.index(slide) + 1}.")
        # print(f"--Prev_file: {prev_file_number}, prev_station_depth: {prev_station_depth}")
        #print(f"--Curr_file: {current_file_number}, curr_station_depth: {current_station_depth}")
        # flag to avoid sample capture image duplication in case multiple text lables exist on a single slide (e.g. "1 x 250cc SPMC")
        first_sample_capture_label = False
        for shape in slide.shapes:
            # print(f"SHAPE: {shape.shape_type}")
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                try:
                    current_file_number = int(re.search("(?<=[Ff]ile )(\d+)(?=[ ]{0,2}:)", shape.text).group())
                    # print(f"---Current file no: {current_file_number}")
                    current_station_depth = float(re.search('(\d{4,5}\.{0,1}\d{0,1})(?=[ ]{0,2}ft)', shape.text).group())
                    # print(f"---Current station depth: {current_station_depth}")

                    if current_file_number != prev_file_number or current_station_depth != prev_station_depth:
                        # print('DEPTH/FILE CHANGE!')
                        for shape in slide.shapes:
                            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                prev_file_number = current_file_number
                                prev_station_depth = current_station_depth
                                shape.Name = f"{current_file_number}_{current_station_depth}_DFA"
                                yield shape
                                continue
                    else:
                        if re.search("\([ ]{0,1}[Cc]omplete", shape.text):
                            # print("This is Integration plot slide")
                            for shape in slide.shapes:
                                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                    shape.Name = f"{current_file_number}_{current_station_depth}_Complete_station"
                                    yield shape
                                    continue
                        elif re.search("[Ss]ampl.+[Cc]aptur", shape.text):
                            sample_capture_counter = 1
                            # check if sample capture plots exist for this file-depth
                            sample_capture_plots = []
                            for file in os.listdir(DIRECTORY):
                                if file.startswith(f"{current_file_number}_{current_station_depth}_Sample_capture_") and file.endswith(".png"):
                                    sample_capture_plots.append(file)
                            # if some sample capture plots already exist for this depth - need to take it into account
                            # and update sample capture counter
                            if len(sample_capture_plots) > 0:
                                # there's at least one sample capture plot existing
                                max_sample_capture_number = 1
                                for sample_capture_plot in sample_capture_plots:
                                    current_sample_capture_plot_number = int(re.search("(?<=_)(\d+)(?=.png)", sample_capture_plot).group())
                                    if max_sample_capture_number < current_sample_capture_plot_number:
                                        max_sample_capture_number = current_sample_capture_plot_number

                                sample_capture_counter = max_sample_capture_number + 1


                            for shape in slide.shapes:
                                if ((shape.shape_type == MSO_SHAPE_TYPE.PICTURE or (shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.placeholder_format.idx == 1))) \
                                        and not first_sample_capture_label:
                                    shape.Name = f"{current_file_number}_{current_station_depth}_Sample_capture_{sample_capture_counter}"
                                    sample_capture_counter += 1
                                    yield shape
                            first_sample_capture_label = True
                            continue

                        elif re.search("[Cc]omposition", shape.text):
                            # print("HC COMP SLIDE!")
                            HC_comp_dict = {
                                'slide_number': prs.slides.index(slide)+1,
                                'file_number': prev_file_number,
                                'station_depth': prev_station_depth
                            }
                            #HC_Composition_slides.append(int(prs.slides.index(slide)+1))
                            HC_Composition_slides.append(HC_comp_dict)

                            for shape in slide.shapes:
                                if (shape.shape_type == MSO_SHAPE_TYPE.PICTURE or (shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.placeholder_format.idx == 1)):
                                    shape.Name = f"{current_file_number}_{current_station_depth}_HCcomp"
                                    yield shape
                                    continue

                        elif re.search("[Cc]ontamination", shape.text):
                            # print("CONTAMINATION SLIDE!")
                            for shape in slide.shapes:
                                if (shape.shape_type == MSO_SHAPE_TYPE.PICTURE or (
                                        shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.placeholder_format.idx == 1)):
                                    shape.Name = f"{current_file_number}_{current_station_depth}_Contamination"
                                    yield shape
                                    continue

                except AttributeError:
                    if re.search("[Cc]omposition", shape.text):
                        # print("HC COMP SLIDE!")
                        HC_comp_dict = {
                            'slide_number': prs.slides.index(slide)+1,
                            'file_number': prev_file_number,
                            'station_depth': prev_station_depth
                        }
                        #HC_Composition_slides.append(int(prs.slides.index(slide)+1))
                        HC_Composition_slides.append(HC_comp_dict)

                        for shape in slide.shapes:
                            if (shape.shape_type == MSO_SHAPE_TYPE.PICTURE or (shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.placeholder_format.idx == 1)):
                                shape.Name = f"{current_file_number}_{current_station_depth}_HCcomp"
                                yield shape
                                continue

                    elif re.search("[Ss]ampling [Ss]ummary", shape.text):
                        Sampling_summary_slides.append(prs.slides.index(slide)+1)

                    elif (re.search("\d{0,2}[ ]{0,2}[Xx]{1}[ ]{0,2}\d{3}[ ]{0,2}[Cc]|[Gg]{2}[ ]{0,2}[\r\n]*\w{4}", shape.text) or \
                        re.search("\d{0,2}[ ]{0,2}[Xx]{1}[ ]{0,2}.+([Gg]|[Gg]al)[ ]{0,2}(SC|MRSC)", shape.text)) and \
                            not first_sample_capture_label:
                        sample_capture_counter = 1
                        # check if sample capture plots exist for this file-depth
                        sample_capture_plots = []
                        for file in os.listdir(DIRECTORY):
                            if file.startswith(
                                    f"{current_file_number}_{current_station_depth}_Sample_capture_") and file.endswith(
                                    ".png"):
                                sample_capture_plots.append(file)
                        # if some sample capture plots already exist for this depth - need to take it into account
                        # and update sample capture counter
                        if len(sample_capture_plots) > 0:
                            # there's at least one sample capture plot existing
                            max_sample_capture_number = 1
                            for sample_capture_plot in sample_capture_plots:
                                current_sample_capture_plot_number = int(
                                    re.search("(?<=_)(\d+)(?=.png)", sample_capture_plot).group())
                                if max_sample_capture_number < current_sample_capture_plot_number:
                                    max_sample_capture_number = current_sample_capture_plot_number

                            sample_capture_counter = max_sample_capture_number + 1

                        for shape in slide.shapes:
                            if (shape.shape_type == MSO_SHAPE_TYPE.PICTURE or (shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.placeholder_format.idx == 1)):
                                # print("SAMPLE CAPTURE PIC DETECTED!")
                                shape.Name = f"{current_file_number}_{current_station_depth}_Sample_capture_{sample_capture_counter}"
                                sample_capture_counter += 1
                                yield shape
                        first_sample_capture_label = True

                    elif re.search("[Cc]ontamination", shape.text):
                        # print("CONTAMINATION SLIDE!")
                        for shape in slide.shapes:
                            if (shape.shape_type == MSO_SHAPE_TYPE.PICTURE or (shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.placeholder_format.idx == 1)):
                                shape.Name = f"{current_file_number}_{current_station_depth}_Contamination"
                                yield shape
                                continue
                    continue

def image_export(FILENAME):
    # export images from ppt
    count = 0
    for picture in iter_shapes(Presentation(FILENAME)):
        count = count + 1
        image = picture.image
        # ---get image "file" contents---
        image_bytes = image.blob
        # ---make up a name for the file, e.g. 'image.jpg'---
        #image_filename = '%s.%s' % (count, image.ext)
        image_filename = f"{picture.Name}.{image.ext}"
        with open(image_filename, 'wb') as f:
            f.write(image_bytes)

        # crop image is height is too large
        im = Image.open(image_filename)
        w, h = im.size
        # image height in inches, assuming dpi = 120
        image_height_inches = im.height / 120
        # print(image_height_inches)
        if image_height_inches > 10.6:
            excess_height_inches = image_height_inches - 10.6
            excess_height_px = int(excess_height_inches * 120)
            # crop excess height from top
            im.crop((0, excess_height_px, w, h)).save(image_filename)

        # print('WIDTH:' + str(im.width))
        # print('HEIGHT:' + str(im.height))
        # print('WIDTH in inches:' + str(im.width / im.info['dpi'][1]))
        # print('HEIGHT in inches: ' + str(im.height / im.info['dpi'][1]))



        # add picture border
        add_border(image_filename, output_image=image_filename, border=1)

    # export sampling summary and HC composition tables from ppt as images
    APPLICATION = win32com.client.gencache.EnsureDispatch('Powerpoint.Application')
    APPLICATION.Visible = True
    PRESENTATION = APPLICATION.Presentations.Open(f'{DIRECTORY}\\{FILENAME}', ReadOnly=False)

    # export Sampling summary as images
    for slide_number in Sampling_summary_slides:
        slide = PRESENTATION.Slides(slide_number)
        for shape in slide.Shapes:
            # (if shape type is table), 19 = JPG image
            if shape.Type == 19:
                shape.Export(f'{DIRECTORY}\\Sampling_summary_{slide_number}.png', 1)

    # export HC composition tables as images
    for slide_dict in HC_Composition_slides:
        slide = PRESENTATION.Slides(slide_dict['slide_number'])
        for shape in slide.Shapes:
            # (if shape type is table), 19 = JPG image
            if shape.Type == 19:
                shape.Export(
                    f'{DIRECTORY}\\{slide_dict["file_number"]}_{slide_dict["station_depth"]}_HCcompTable.png', 1)

    # close PowerPoint app
    APPLICATION.Quit()
