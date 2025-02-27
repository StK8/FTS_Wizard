import collections.abc
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re
import win32com.client
import os

# tool type for each particular station  - by default 'MDT'
CURRENT_TOOL_TYPE = 'MDT'
# tool types used throughout the job - can be either MDT or ORA for now
JOB_TOOL_TYPES = set()

# auxiliary function that returns valid xml chars
def valid_xml_char_ordinal(c):
    codepoint = ord(c)
    # conditions ordered by presumed frequency
    return (
        0x20 <= codepoint <= 0xD7FF or
        codepoint in (0x9, 0xA, 0xD) or
        0xE000 <= codepoint <= 0xFFFD or
        0x10000 <= codepoint <= 0x10FFFF
        )

# parsing ppt sampling summary tables
def table_parse(FILENAME):
    prs = Presentation(FILENAME)
    DFA_stations_raw = []
    for slide in prs.slides:
        shape_table = None
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                if re.search("[Ss]ampling [Ss]ummary", shape.text):
                    for shape in slide.shapes:
                        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                            shape_table = shape
                            table = shape_table.table
                            for row in table.rows:
                                DFA_station_raw = {}
                                if not re.search("[Zz]one|Run[ ]{0,2}\d{1}[ ]{0,2}\w{0,4}", row.cells[0].text):
                                    DFA_station_raw['formation'] = row.cells[0].text
                                    DFA_station_raw['file_number'] = row.cells[1].text
                                    DFA_station_raw['depth'] = str(float(row.cells[2].text))
                                    DFA_station_raw['time'] = row.cells[3].text
                                    DFA_station_raw['volume'] = row.cells[4].text
                                    DFA_station_raw['max_dd'] = row.cells[5].text
                                    DFA_station_raw['samples_number'] = row.cells[6].text
                                    DFA_station_raw['probe'] = row.cells[7].text
                                    DFA_station_raw['mobility'] = row.cells[8].text
                                    DFA_station_raw['pressure'] = row.cells[9].text
                                    DFA_station_raw['temperature'] = row.cells[10].text
                                    DFA_station_raw['observation'] = row.cells[11].text
                                    DFA_station_raw['bottles_observation'] = row.cells[12].text
                                    DFA_station_raw['comments'] = row.cells[13].text
                                    DFA_station_raw['tool_type'] = CURRENT_TOOL_TYPE
                                    DFA_stations_raw.append(DFA_station_raw)
                                elif re.search("[Rr][Uu][Nn][ ]{0,3}\d[ ]{0,3}", row.cells[0].text):
                                    if re.search("[Mm][Dd][Tt]", row.cells[0].text):
                                        CURRENT_TOOL_TYPE = 'MDT'
                                        JOB_TOOL_TYPES.add(CURRENT_TOOL_TYPE)
                                    elif re.search("[Oo][Rr][Aa]", row.cells[0].text):
                                        CURRENT_TOOL_TYPE = 'ORA'
                                        JOB_TOOL_TYPES.add(CURRENT_TOOL_TYPE)
    # print(DFA_stations_raw)

    return DFA_stations_raw

# processing parsed ppt sampling summary tables' data
def table_processing(DFA_stations_raw):
    DFA_stations_processed = []

    for DFA_station_raw in DFA_stations_raw:
        DFA_station_processed = {}

        # checking if formation name exists
        if re.search("\w", DFA_station_raw['formation']):
            DFA_station_processed['formation'] = DFA_station_raw['formation']
        else:
            DFA_station_processed['formation'] = '-'

        DFA_station_processed['file_number'] = DFA_station_raw['file_number']
        DFA_station_processed['depth'] = DFA_station_raw['depth']

        # checking if PO time exists
        if re.search("\d", DFA_station_raw['time']):
            DFA_station_processed['time'] = DFA_station_raw['time']
        else:
            DFA_station_processed['time'] = '-'

        # checking if PO volume exists
        if re.search("\d", DFA_station_raw['volume']):
            DFA_station_processed['volume'] = DFA_station_raw['volume']
        else:
            DFA_station_processed['volume'] = '-'

        # checking if max drawdown values exist
        # if both max DD PO and max DD samplng values exist
        if re.search("\/", DFA_station_raw['max_dd']):
            max_dd = DFA_station_raw['max_dd'].strip('~ ')
            DFA_station_processed['max_dd_po'] = re.search('^(\d+)', max_dd).group()
            DFA_station_processed['max_dd_sampling'] = re.search('(\d+)$', max_dd).group()
        # else if only max DD PO value exists
        elif re.search("\d+", DFA_station_raw['max_dd']):
            max_dd_po = DFA_station_raw['max_dd'].strip('~ ')
            DFA_station_processed['max_dd_po'] = re.search('\d+', max_dd_po).group()
            DFA_station_processed['max_dd_sampling'] = '-'
        # if no DD values are available
        else:
            DFA_station_processed['max_dd_po'] = '-'
            DFA_station_processed['max_dd_sampling'] = '-'

        # tool type used at particular station
        DFA_station_processed['tool_type'] = DFA_station_raw['tool_type']

        # checking number of samples captured
        # if no samples were captured
        if re.search('DFA|-', DFA_station_raw['samples_number']):
            DFA_station_processed['station_type'] = 'DFA'
            DFA_station_processed['mpsr_number'] = 0
            DFA_station_processed['spmc_number'] = 0
            # 1st element of the list - 1 gal SC, 2nd list - 2.75 gal SC, 3rd list - 6 gal SC
            DFA_station_processed['mrsc_number'] = [[0], [0], [0]]
            DFA_station_processed['fnlt_number'] = 0
            DFA_station_processed['fnst_number'] = 0

        # if samples were captured
        elif re.search('MPSR|SPMC|FNLT|FNST|SC', DFA_station_raw['samples_number']):
            # put initial samples numbers as 0
            DFA_station_processed['mpsr_number'] = 0
            DFA_station_processed['spmc_number'] = 0
            # 1st element of the list - 1 gal SC, 2nd list - 2.75 gal SC, 3rd list - 6 gal SC
            DFA_station_processed['mrsc_number'] = [[0], [0], [0]]
            DFA_station_processed['fnlt_number'] = 0
            DFA_station_processed['fnst_number'] = 0
            DFA_station_processed['station_type'] = 'Sampling'
            # for number of samples format in ppt like: 2x450 cc MPSR
            if re.search('450', DFA_station_raw['samples_number']):
                # find all occurences
                mpsr_occurences = re.findall('(\d+)(?=[ ]{0,2}x[ ]{0,2}450)', DFA_station_raw['samples_number'])
                for mpsr_occurence in mpsr_occurences:
                    DFA_station_processed['mpsr_number'] += int(mpsr_occurence)
                #DFA_station_processed['mpsr_number'] = int(re.search('(\d+)(?=[ ]{0,2}x[ ]{0,2}450)', DFA_station_raw['samples_number']).group())
            # for number of samples format in ppt like: 2xMPSR
            elif re.search('MPSR', DFA_station_raw['samples_number']):
                # find all occurences
                mpsr_occurences = re.findall('(\d+)(?=[ ]{0,2}x[ ]{0,2}MPSR)', DFA_station_raw['samples_number'])
                for mpsr_occurence in mpsr_occurences:
                    DFA_station_processed['mpsr_number'] += int(mpsr_occurence)
                #DFA_station_processed['mpsr_number'] = int(re.search('(\d+)(?=[ ]{0,2}x[ ]{0,2}MPSR)', DFA_station_raw['samples_number']).group())
            # for number of samples format in ppt like: 2x250 cc SPMC
            if re.search('250', DFA_station_raw['samples_number']):
                # find all occurences
                spmc_occurences = re.findall('(\d+)(?=[ ]{0,2}x[ ]{0,2}250)', DFA_station_raw['samples_number'])
                for spmc_occurence in spmc_occurences:
                    DFA_station_processed['spmc_number'] += int(spmc_occurence)
                #DFA_station_processed['spmc_number'] = int(re.search('(\d+)(?=[ ]{0,2}x[ ]{0,2}250)', DFA_station_raw['samples_number']).group())
            # for number of samples format in ppt like: 2xSPMC
            elif re.search('SPMC', DFA_station_raw['samples_number']):
                # find all occurences
                spmc_occurences = re.findall('(\d+)(?=[ ]{0,2}x[ ]{0,2}SPMC)', DFA_station_raw['samples_number'])
                for spmc_occurence in spmc_occurences:
                    DFA_station_processed['spmc_number'] += int(spmc_occurence)
                #DFA_station_processed['spmc_number'] = int(re.search('(\d+)(?=[ ]{0,2}x[ ]{0,2}SPMC)', DFA_station_raw['samples_number']).group())
            # for number of samples format in ppt like: 2x675 cc FNLT
            if re.search('675', DFA_station_raw['samples_number']):
                # find all occurences
                fnlt_occurences = re.findall('(\d+)(?=[ ]{0,2}x[ ]{0,2}675)', DFA_station_raw['samples_number'])
                for fnlt_occurence in fnlt_occurences:
                    DFA_station_processed['fnlt_number'] += int(fnlt_occurence)
                #DFA_station_processed['fnlt_number'] = int(re.search('(\d+)(?=[ ]{0,2}x[ ]{0,2}675)', DFA_station_raw['samples_number']).group())
            # for number of samples format in ppt like: 2xFNLT
            elif re.search('FNLT', DFA_station_raw['samples_number']):
                # find all occurences
                fnlt_occurences = re.findall('(\d+)(?=[ ]{0,2}x[ ]{0,2}FNLT)', DFA_station_raw['samples_number'])
                for fnlt_occurence in fnlt_occurences:
                    DFA_station_processed['fnlt_number'] += int(fnlt_occurence)
                #DFA_station_processed['fnlt_number'] = int(re.search('(\d+)(?=[ ]{0,2}x[ ]{0,2}FNLT)', DFA_station_raw['samples_number']).group())
            # for number of samples format in ppt like: 2x400 cc FNST
            if re.search('400', DFA_station_raw['samples_number']):
                # find all occurences
                fnst_occurences = re.findall('(\d+)(?=[ ]{0,2}x[ ]{0,2}400)', DFA_station_raw['samples_number'])
                for fnst_occurence in fnst_occurences:
                    DFA_station_processed['fnst_number'] += int(fnst_occurence)
                #DFA_station_processed['fnst_number'] = int(re.search('(\d+)(?=[ ]{0,2}x[ ]{0,2}400)', DFA_station_raw['samples_number']).group())
            # for number of samples format in ppt like: 2xFNST
            elif re.search('FNST', DFA_station_raw['samples_number']):
                # find all occurences
                fnst_occurences = re.findall('(\d+)(?=[ ]{0,2}x[ ]{0,2}FNST)', DFA_station_raw['samples_number'])
                for fnst_occurence in fnst_occurences:
                    DFA_station_processed['fnst_number'] += int(fnst_occurence)
                #DFA_station_processed['fnst_number'] = int(re.search('(\d+)(?=[ ]{0,2}x[ ]{0,2}FNST)', DFA_station_raw['samples_number']).group())
            # for 1 gal SC chambers - if volume in gallons
            if re.search('1(\.){0,1}[0]{0,1}[ ]{0,2}[Gg].+SC', DFA_station_raw['samples_number']):
                DFA_station_processed['mrsc_number'][0][0] = int(re.search('(\d+)(?=[ ]{0,2}x[ ]{0,2}1[\.]{0,1}[0]{0,1}[ ]{0,1}[Gg].+)',DFA_station_raw['samples_number']).group())
            # for 1 gal SC chambers - if volume in liters or cc
            elif re.search('3.+[Ll|cc].+SC', DFA_station_raw['samples_number']):
                DFA_station_processed['mrsc_number'][0][0] = int(re.search('(\d+)(?=[ ]{0,2}x[ ]{0,2}3)', DFA_station_raw['samples_number']).group())
            # for 2.75 gal SC chambers - if volume in gallons
            if re.search('2\.7[5]{0,1}[ ]{0,2}[Gg].+SC', DFA_station_raw['samples_number']):
                DFA_station_processed['mrsc_number'][1][0] = int(re.search('(\d+)(?=[ ]{0,2}x[ ]{0,2}2.7)', DFA_station_raw['samples_number']).group())
            # for 2.75 gal SC chambers - if volume in liters
            elif re.search('10.+[Ll].+SC', DFA_station_raw['samples_number']):
                DFA_station_processed['mrsc_number'][1][0] = int(re.search('(\d+)(?=[ ]{0,2}x[ ]{0,2}10)', DFA_station_raw['samples_number']).group())
            # for 6 gal SC chambers - if volume in gallons
            if re.search('6(\.){0,1}[0]{0,1}[ ]{0,2}[Gg].+SC', DFA_station_raw['samples_number']):
                DFA_station_processed['mrsc_number'][2][0] = int(re.search('(\d+)(?=[ ]{0,2}x[ ]{0,2}6[\.]{0,1}[0]{0,1}[ ]{0,1}[Gg].+)', DFA_station_raw['samples_number']).group())
            # for 6 gal SC chambers - if volume in liters
            elif re.search('22.+[Ll].+SC', DFA_station_raw['samples_number']):
                DFA_station_processed['mrsc_number'][2][0] = int(re.search('(\d+)(?=[ ]{0,2}x[ ]{0,2}22)', DFA_station_raw['samples_number']).group())


        # if probe XLD or LD - add word 'probe'
        if (re.search('XLD', DFA_station_raw['probe']) or re.search('LD', DFA_station_raw['probe'])) and not \
            re.search('[Pp]robe]', DFA_station_raw['probe']):
            DFA_station_processed['probe'] = DFA_station_raw['probe'] + ' probe'
        else:
            DFA_station_processed['probe'] = DFA_station_raw['probe']

        # checking if mobility exists
        if re.search("\d", DFA_station_raw['mobility']):
            mobility = DFA_station_raw['mobility'].strip('~ ')
            DFA_station_processed['mobility'] = mobility
        else:
            DFA_station_processed['mobility'] = '-'

        # checking if formation pressure exists
        if re.search("\d", DFA_station_raw['pressure']):
            DFA_station_processed['pressure'] = re.search('\d{3,5}\.{0,1}\d{0,2}', DFA_station_raw['pressure']).group()
        else:
            DFA_station_processed['pressure'] = '-'

        # checking if temperature exists
        if re.search("\d", DFA_station_raw['temperature']):
            DFA_station_processed['temperature'] = re.search('\d{2,3}\.{0,1}\d{0,2}', DFA_station_raw['temperature']).group()
        else:
            DFA_station_processed['temperature'] = '-'

        # fluid observed
        observation = DFA_station_raw['observation'].strip('\n')
        observation = observation.replace('\n', ' ')
        observation = observation.replace('  ', ' ')
        if re.search('[Dd]ensity', observation):
            observation = re.search('(.+)(?=[Dd]ensity)', observation).group()
        if re.search('[Gg]OR', observation):
            observation = re.search('(.+)(?=[Gg]OR)', observation).group()
        if re.search('[Rr]esistivity', observation):
            observation = re.search('(.+)(?=[Rr]esistivity)', observation).group()
        if re.search('[Ss]alinity', observation):
            observation = re.search('(.+)(?=[Ss]alinity)', observation).group()
        observation = observation.strip('*^ ')
        DFA_station_processed['observation'] = observation

        # bottle types and serial numbers
        bottles = []
        mpsr_bottles = []
        spmc_bottles = []
        fnlt_bottles = []
        fnst_bottles = []
        mrsc_1gal_bottles = []
        mrsc_2gal_bottles = []
        mrsc_6gal_bottles = []

        # check MPSR serial numbers
        if DFA_station_processed['mpsr_number'] > 0:
            match = re.findall('(?<=MPSR)(.+)(?=\*|\^|\n|\s)', DFA_station_raw['bottles_observation'])
            for bottle in match:
                bottles.append(f'MPSR {bottle.strip("*^# ")}')
                mpsr_bottles.append(f'MPSR {bottle.strip("*^# ")}')
        # check SPMC serial numbers
        if DFA_station_processed['spmc_number'] > 0:
            match = re.findall('(?<=SPMC)(.+)(?=\*|\^|\n|\s)', DFA_station_raw['bottles_observation'])
            for bottle in match:
                bottles.append(f'SPMC {bottle.strip("*^# ")}')
                spmc_bottles.append(f'SPMC {bottle.strip("*^# ")}')
        # check 1 gal MRSC serial numbers
        if DFA_station_processed['mrsc_number'][0][0] > 0 and DFA_station_processed['mrsc_number'][1][0] == 0 and \
                DFA_station_processed['mrsc_number'][2][0] == 0:
            match = re.findall('(?<=SC)(.+)(?=\*|\^|\n|\s)', DFA_station_raw['bottles_observation'])
            for bottle in match:
                bottles.append(f'MRSC {bottle.strip("*^# ")}')
                mrsc_1gal_bottles.append(f'MRSC {bottle.strip("*^# ")}')
        # check 2.75 gal MRSC serial numbers
        if DFA_station_processed['mrsc_number'][1][0] > 0 and DFA_station_processed['mrsc_number'][0][0] == 0 and \
                DFA_station_processed['mrsc_number'][2][0] == 0:
            match = re.findall('(?<=SC)(.+)(?=\*|\^|\n|\s)', DFA_station_raw['bottles_observation'])
            for bottle in match:
                bottles.append(f'MRSC {bottle.strip("*^# ")}')
                mrsc_2gal_bottles.append(f'MRSC {bottle.strip("*^# ")}')
        # check 6 gal MRSC serial numbers
        if DFA_station_processed['mrsc_number'][2][0] > 0 and DFA_station_processed['mrsc_number'][0][0] == 0 and \
                DFA_station_processed['mrsc_number'][1][0] == 0:
            match = re.findall('(?<=SC)(.+)(?=\*|\^|\n|\s)', DFA_station_raw['bottles_observation'])
            for bottle in match:
                bottles.append(f'MRSC {bottle.strip("*^# ")}')
                mrsc_6gal_bottles.append(f'MRSC {bottle.strip("*^# ")}')
        # check FNLT serial numbers
        if DFA_station_processed['fnlt_number'] > 0:
            match = re.findall('(?<=FNLT)(.+)(?=\*|\^|\n|\s)', DFA_station_raw['bottles_observation'])
            for bottle in match:
                bottles.append(f'FNLT {bottle.strip("*^# ")}')
                fnlt_bottles.append(f'FNLT {bottle.strip("*^# ")}')
        # check FNST serial numbers
        if DFA_station_processed['fnst_number'] > 0:
            match = re.findall('(?<=FNST)(.+)(?=\*|\^|\n|\s)', DFA_station_raw['bottles_observation'])
            for bottle in match:
                bottles.append(f'FNST {bottle.strip("*^# ")}')
                fnst_bottles.append(f'FNST {bottle.strip("*^# ")}')
        # checking that all bottles serial numbers were parsed properly
        if (DFA_station_processed['mpsr_number'] + DFA_station_processed['spmc_number'] +
                DFA_station_processed['fnlt_number'] + DFA_station_processed['fnst_number'] +
                DFA_station_processed['mrsc_number'][0][0] + \
                DFA_station_processed['mrsc_number'][1][0] + \
                DFA_station_processed['mrsc_number'][2][0]) != len(bottles):
            bottles = ['N/A']
        DFA_station_processed['bottles_observation'] = bottles
        DFA_station_processed['mpsr_bottles'] = mpsr_bottles
        DFA_station_processed['spmc_bottles'] = spmc_bottles
        DFA_station_processed['mrsc_1gal_bottles'] = mrsc_1gal_bottles
        DFA_station_processed['mrsc_2gal_bottles'] = mrsc_2gal_bottles
        DFA_station_processed['mrsc_6gal_bottles'] = mrsc_6gal_bottles
        DFA_station_processed['fnlt_bottles'] = fnlt_bottles
        DFA_station_processed['fnst_bottles'] = fnst_bottles

        # checking if comments exist
        if re.search("\w", DFA_station_raw['comments']):
            # process the text
            comments = DFA_station_raw['comments']
            comments = comments.replace('\n', ' ')
            comments = comments.strip('*^ ')
            comments = f'{comments}.'
            DFA_station_processed['comments'] = comments
        else:
            DFA_station_processed['comments'] = '-'

        DFA_stations_processed.append(DFA_station_processed)

    # for DFA_station_processed in DFA_stations_processed:
    #     print(DFA_station_processed)
    return DFA_stations_processed


# generate text for each DFA/Sampling station
def generate_stations_text(DFA_stations_processed):

    DFA_stations_text = []
    for DFA_station in DFA_stations_processed:
        # if no pumping took place for whatever reason
        if (DFA_station['volume'] == '-' or DFA_station['time'] == '-'):
            text = 'DFA station was done at ' + DFA_station['depth'] + ' ft MD with ' \
                   + DFA_station['probe'] + '. ' + DFA_station['comments']
        else:
            text = DFA_station['station_type'] + ' station was done at ' + DFA_station['depth'] \
                   + ' ft MD with ' + DFA_station['probe'] + '. A total of ' + DFA_station['volume'] \
                   + ' liters of fluid was pumped out for ' + DFA_station['time'] + ' hrs. Fluid was identified as ' \
                   + DFA_station['observation'] + '. Maximum pressure drawdown during pumpout was ' \
                   + DFA_station['max_dd_po'] + ' psi '

            # if sampling station  -add max dd during sampling
            if DFA_station['station_type'] == 'Sampling':
                text += 'and maximum pressure drawdown during sampling was ' + DFA_station['max_dd_sampling'] + ' psi '

            text += 'with pumpout mobility of ' + DFA_station['mobility'] + ' mD/cP. '

            # if sampling station  -add bottles captured text
            if DFA_station['station_type'] == 'Sampling':
                if (DFA_station['mpsr_number'] + DFA_station['spmc_number'] + int(DFA_station['mrsc_number'][0][0]) +
                    DFA_station['mrsc_number'][1][0] + DFA_station['mrsc_number'][2][0]) == 1 or \
                        (DFA_station['fnlt_number'] + DFA_station['fnst_number']) == 1:
                    text += 'The following sample was captured during the station: '
                elif (DFA_station['mpsr_number'] + DFA_station['spmc_number'] + DFA_station['mrsc_number'][0][0] +
                    DFA_station['mrsc_number'][1][0] + DFA_station['mrsc_number'][2][0]) > 1 or \
                        (DFA_station['fnlt_number'] + DFA_station['fnst_number']) > 1:
                    text += 'The following samples were captured during the station: '

                # if both MPSR and SPMC bottles were captured
                if DFA_station['mpsr_number'] > 0 and DFA_station['spmc_number'] > 0:
                    # text related to MPSR bottles
                    mpsr_text = str(DFA_station['mpsr_number']) + 'x450cc ('
                    for bottle in DFA_station['mpsr_bottles']:
                        mpsr_text += bottle + ', '
                    mpsr_text = mpsr_text.strip(', ')
                    mpsr_text += ')'

                    # text related to SPMC bottles
                    spmc_text = str(DFA_station['spmc_number']) + 'x250cc ('
                    for bottle in DFA_station['spmc_bottles']:
                        spmc_text += bottle + ', '
                    spmc_text = spmc_text.strip(', ')
                    spmc_text += ')'

                    text += mpsr_text + ' and ' + spmc_text + '. '

                # if only MPSR bottle(s) were captured
                elif DFA_station['mpsr_number'] > 0:
                    # text related to MPSR bottles
                    mpsr_text = str(DFA_station['mpsr_number']) + 'x450cc ('
                    for bottle in DFA_station['mpsr_bottles']:
                        mpsr_text += bottle + ', '
                    mpsr_text = mpsr_text.strip(', ')
                    mpsr_text += ')'

                    text += mpsr_text + '. '

                # if only SPMC bottle(s) were captured
                elif DFA_station['spmc_number'] > 0:
                    # text related to SPMC bottles
                    spmc_text =str(DFA_station['spmc_number']) + 'x250cc ('
                    for bottle in DFA_station['spmc_bottles']:
                        spmc_text += bottle + ', '
                    spmc_text = spmc_text.strip(', ')
                    spmc_text += ')'

                    text += spmc_text + '. '

                # if only 1 gal MRSC chamber was filled
                elif DFA_station['mrsc_number'][0][0] > 0:
                    # text related to 1.0 gal MRSC chambers
                    mrsc_text =str(DFA_station['mrsc_number'][0][0]) + 'x1.0 gal ('
                    for bottle in DFA_station['mrsc_1gal_bottles']:
                        mrsc_text += bottle + ', '
                    mrsc_text = mrsc_text.strip(', ')
                    mrsc_text += ')'

                    text += mrsc_text + '. '

                # if only 2.75 gal MRSC chamber was filled
                elif DFA_station['mrsc_number'][1][0] > 0:
                    # text related to 2.75 gal MRSC chambers
                    mrsc_text =str(DFA_station['mrsc_number'][1][0]) + 'x2.75 gal ('
                    for bottle in DFA_station['mrsc_2gal_bottles']:
                        mrsc_text += bottle + ', '
                    mrsc_text = mrsc_text.strip(', ')
                    mrsc_text += ')'

                    text += mrsc_text + '. '

                # if only 6 gal MRSC chamber was filled
                elif DFA_station['mrsc_number'][2][0] > 0:
                    # text related to 6.0 gal MRSC chambers
                    mrsc_text =str(DFA_station['mrsc_number'][2][0]) + 'x6.0 gal ('
                    for bottle in DFA_station['mrsc_6gal_bottles']:
                        mrsc_text += bottle + ', '
                    mrsc_text = mrsc_text.strip(', ')
                    mrsc_text += ')'

                    text += mrsc_text + '. '

                # if both FNLT and FNST bottles were captured
                elif DFA_station['fnlt_number'] > 0 and DFA_station['fnst_number'] > 0:
                    # text related to FNLT bottles
                    fnlt_text = str(DFA_station['fnlt_number']) + 'x675cc ('
                    for bottle in DFA_station['fnlt_bottles']:
                        fnlt_text += bottle + ', '
                    fnlt_text = fnlt_text.strip(', ')
                    fnlt_text += ')'

                    # text related to FNST bottles
                    fnst_text = str(DFA_station['fnst_number']) + 'x400cc ('
                    for bottle in DFA_station['fnst_bottles']:
                        fnst_text += bottle + ', '
                    fnst_text = fnst_text.strip(', ')
                    fnst_text += ')'

                    text += fnlt_text + ' and ' + fnst_text + '. '

                # if only FNLT bottle(s) were captured
                elif DFA_station['fnlt_number'] > 0:
                    # text related to FNLT bottles
                    fnlt_text = str(DFA_station['fnlt_number']) + 'x675cc ('
                    for bottle in DFA_station['fnlt_bottles']:
                        fnlt_text += bottle + ', '
                    fnlt_text = fnlt_text.strip(', ')
                    fnlt_text += ')'

                    text += fnlt_text + '. '

                # if only FNST bottle(s) were captured
                elif DFA_station['fnst_number'] > 0:
                    # text related to FNST bottles
                    fnst_text = str(DFA_station['fnst_number']) + 'x400cc ('
                    for bottle in DFA_station['fnst_bottles']:
                        fnst_text += bottle + ', '
                    fnst_text = fnst_text.strip(', ')
                    fnst_text += ')'

                    text += fnst_text + '. '

            text += DFA_station['comments']
        # print(text)
        station = {}
        station['file_number'] = DFA_station['file_number']
        station['depth'] = DFA_station['depth']
        station['station_type'] = DFA_station['station_type']
        station['tool_type'] = DFA_station['tool_type']
        station['time'] = DFA_station['time']
        station['volume'] = DFA_station['volume']
        # remove invalid xml chars
        cleaned_string = ''.join(c for c in text if valid_xml_char_ordinal(c))
        cleaned_string = cleaned_string.replace('\n', ' ')
        station['text'] = cleaned_string

        DFA_stations_text.append(station)
    return DFA_stations_text

# generate text for 'Key observations' section of the report
def generate_stations_summary_text(DFA_stations_processed):
    DFA_stations_summary_text = []
    for DFA_station in DFA_stations_processed:
        text = ''
        # if no pumping took place for whatever reason
        if (DFA_station['volume'] == '-' or DFA_station['time'] == '-'):
            text += 'DFA station at ' + DFA_station['depth'] + ' ft MD was conducted ' \
                + DFA_station['comments']
        elif DFA_station['station_type'] == 'DFA':
            text += 'DFA station at ' + DFA_station['depth'] + ' ft MD was conducted ' \
                + 'where ' + DFA_station['observation'] + ' was observed'
        elif DFA_station['station_type'] == 'Sampling':
            text += 'Sampling station at ' + DFA_station['depth'] + ' ft MD was conducted ' \
                + 'where ' + DFA_station['observation'] + ' was observed and sampled'
        DFA_stations_summary_text.append(text)
        # print(text)
    return DFA_stations_summary_text
